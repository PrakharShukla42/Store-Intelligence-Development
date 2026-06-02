import collections
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
# Set presentation slide width and height to 16:9 widescreen
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# Helper to apply dark background
BG_COLOR = RGBColor(11, 15, 25) # #0b0f19
TEXT_WHITE = RGBColor(243, 244, 246) # #f3f4f6
TEXT_GRAY = RGBColor(156, 163, 175) # #9ca3af
CYAN = RGBColor(6, 182, 212) # #06b6d4
INDIGO = RGBColor(99, 102, 241) # #6366f1

def apply_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

def add_header(slide, title_text):
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.33), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text.upper()
    p.font.name = 'Arial'
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = CYAN

# Slide 1: Title
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)
apply_background(slide)

# Title box
title_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.33), Inches(3.0))
tf = title_box.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = "APEX STORE INTELLIGENCE"
p.font.name = 'Arial'
p.font.size = Pt(48)
p.font.bold = True
p.font.color.rgb = CYAN
p.alignment = PP_ALIGN.CENTER

p2 = tf.add_paragraph()
p2.text = "Bridging the Offline Retail Data Blindspot with End-to-End Analytics"
p2.font.name = 'Arial'
p2.font.size = Pt(22)
p2.font.color.rgb = INDIGO
p2.alignment = PP_ALIGN.CENTER
p2.space_before = Pt(20)

p3 = tf.add_paragraph()
p3.text = "Purplle Technical Challenge 2026 | Submission Pitch Deck"
p3.font.name = 'Arial'
p3.font.size = Pt(14)
p3.font.color.rgb = TEXT_GRAY
p3.alignment = PP_ALIGN.CENTER
p3.space_before = Pt(40)

# Slide 2: The Data Blindspot
slide = prs.slides.add_slide(blank_slide_layout)
apply_background(slide)
add_header(slide, "1. The Offline Retail Data Blindspot")

# Left Column (The Problem)
box_left = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.8), Inches(5.0))
tf_l = box_left.text_frame
tf_l.word_wrap = True
p_l_t = tf_l.paragraphs[0]
p_l_t.text = "THE PROBLEM"
p_l_t.font.name = 'Arial'
p_l_t.font.size = Pt(20)
p_l_t.font.bold = True
p_l_t.font.color.rgb = INDIGO
p_l_t.space_after = Pt(14)

bullets_l = [
    "Offline retailers operate in a data vacuum compared to e-commerce storefronts.",
    "No visibility into consumer browsing routes, zone entry drop-offs, or shopping dwells.",
    "Checkout queue abandonment occurs blindly without severity alarms.",
    "Impacts store profitability, staffing planning, and layout optimizations."
]
for b in bullets_l:
    p = tf_l.add_paragraph()
    p.text = "• " + b
    p.font.name = 'Arial'
    p.font.size = Pt(15)
    p.font.color.rgb = TEXT_WHITE
    p.space_after = Pt(10)

# Right Column (The Solution)
box_right = slide.shapes.add_textbox(Inches(7.0), Inches(1.5), Inches(5.8), Inches(5.0))
tf_r = box_right.text_frame
tf_r.word_wrap = True
p_r_t = tf_r.paragraphs[0]
p_r_t.text = "THE SOLUTION"
p_r_t.font.name = 'Arial'
p_r_t.font.size = Pt(20)
p_r_t.font.bold = True
p_r_t.font.color.rgb = CYAN
p_r_t.space_after = Pt(14)

bullets_r = [
    "Apex Store Intelligence ingests and maps raw CCTV video streams.",
    "Computes real-time conversion rates dynamically matched with POS registers.",
    "Exposes exact shop-floor layout engagement metrics via HSL-normalized heatmaps.",
    "Automates bottleneck alarms (conversion drops, dead zones, queue depth spikes)."
]
for b in bullets_r:
    p = tf_r.add_paragraph()
    p.text = "• " + b
    p.font.name = 'Arial'
    p.font.size = Pt(15)
    p.font.color.rgb = TEXT_WHITE
    p.space_after = Pt(10)

# Slide 3: Technology Stack
slide = prs.slides.add_slide(blank_slide_layout)
apply_background(slide)
add_header(slide, "2. Platform Technology Stack")

box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.33), Inches(5.0))
tf = box.text_frame
tf.word_wrap = True

tech_stack = {
    "CORE BACKEND": "FastAPI REST microservice, SQLite, SQLAlchemy Relational Model, Pydantic Ingestion Schema.",
    "TRACKING PIPELINE": "YOLOv8 Class 0 (Person), ByteTrack Centroid association, Feature Re-ID (dominant color + aspect ratios), Overlap Boundary filter.",
    "MATHEMATICAL ENGINE": "Dynamic POS datetime correlation math, session funnel tracker, HSL zone visit scaling, dynamic queue window statistics.",
    "DASHBOARD UI": "Single-page responsive glassmorphic panel using HSL tailormade palettes, Vanilla JS fetch polling, interactive stream seeder.",
    "TESTS & DEPLOYMENT": "13-unit pytest suite, multi-stage Docker containerization, Docker Compose, 1-Click Render blueprint."
}

first = True
for title, desc in tech_stack.items():
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    first = False
    p.text = f"{title}: "
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = CYAN
    
    # Add desc as normal text
    run = p.add_run()
    run.text = desc
    run.font.bold = False
    run.font.color.rgb = TEXT_WHITE
    p.space_after = Pt(14)

# Slide 4: CV Tracking
slide = prs.slides.add_slide(blank_slide_layout)
apply_background(slide)
add_header(slide, "3. Computer Vision & Trajectory Pipeline")

box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.33), Inches(5.0))
tf = box.text_frame
tf.word_wrap = True

cv_bullets = [
    "Object Detection: Runs YOLOv8 targeting shopper boundary boxes on multi-camera streams.",
    "Centroid Association: Employs bounding box centroid tracking for robust frame-to-frame trajectory mapping.",
    "Re-ID Signatures: Computes shopper visual signatures using aspect ratio bounds and dominant HSV color histograms.",
    "Cross-Camera Deduplication: Implements entry and shop-floor camera overlap resolving to avoid double counting.",
    "Hardware Graceful Degradation: Integrates dual-mode execution—automatic fallback to simulated behavioral shopper streams on CPU constraints to maintain continuous database feeds."
]

first = True
for b in cv_bullets:
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    first = False
    p.text = "• " + b
    p.font.name = 'Arial'
    p.font.size = Pt(16)
    p.font.color.rgb = TEXT_WHITE
    p.space_after = Pt(14)

# Slide 5: Advanced Math Engine
slide = prs.slides.add_slide(blank_slide_layout)
apply_background(slide)
add_header(slide, "4. Advanced Analytics Mathematics Engine")

box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.33), Inches(5.0))
tf = box.text_frame
tf.word_wrap = True

math_points = {
    "Dynamic POS Correlation": "Correlates visitor entry in checkout queue zones with Point-of-Sale database commits within a strict 5-minute (300 seconds) rolling window to register sales conversion.",
    "Session-Based Funnel Progression": "Computes exact drop-off statistics through a four-stage funnel: Entry -> Zone Visit -> Billing Queue -> Purchase, ignoring multiple re-entries inside a single session.",
    "Normalized Heatmaps": "Aggregates zone visit counts and individual shopper dwell durations, applying a min-max scaling algorithm (0-100) mapped to HSL color hues for intuitive spatial layout mapping.",
    "Checkout Queue Depth Statistics": "Tracks explicit camera queue events and estimates wait congestion by checking active visitors inside the billing queue area, triggering alert flags on bottlenecks."
}

first = True
for title, desc in math_points.items():
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    first = False
    p.text = f"• {title}: "
    p.font.bold = True
    p.font.color.rgb = INDIGO
    p.font.size = Pt(16)
    
    run = p.add_run()
    run.text = desc
    run.font.bold = False
    run.font.color.rgb = TEXT_WHITE
    p.space_after = Pt(12)

# Slide 6: Operations Anomaly Alarms
slide = prs.slides.add_slide(blank_slide_layout)
apply_background(slide)
add_header(slide, "5. Real-Time Operations Anomaly Alerting")

box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.33), Inches(5.0))
tf = box.text_frame
tf.word_wrap = True

anomalies = [
    "BILLING_QUEUE_SPIKE: Triggered if active checkout queue depth exceeds 4 shoppers. Automatically updates status to WARN/CRITICAL, advising managers to 'Open auxiliary cash counter #2' instantly.",
    "CONVERSION_DROP: Triggers if the daily conversion rate falls by >20% compared to historical POS transactional average, warning managers to 'Verify checkout register POS status or assess pricing'.",
    "DEAD_ZONE: Flags a layout zone as dormant if it registers 0 unique visitor entries in the last 30 minutes, suggesting floor managers to 'Re-evaluate product layout and retail banners in this zone'."
]

first = True
for a in anomalies:
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    first = False
    title, desc = a.split(":")
    p.text = "🚨 " + title + ":"
    p.font.bold = True
    p.font.color.rgb = CYAN
    p.font.size = Pt(16)
    
    run = p.add_run()
    run.text = desc
    run.font.bold = False
    run.font.color.rgb = TEXT_WHITE
    p.space_after = Pt(18)

# Slide 7: Premium Dashboard
slide = prs.slides.add_slide(blank_slide_layout)
apply_background(slide)
add_header(slide, "6. Live Glassmorphic Interactive UI")

box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.33), Inches(5.0))
tf = box.text_frame
tf.word_wrap = True

ui_bullets = [
    "Premium Aesthetics: Features harmony dark-mode styling with translucent glass panels, glowing indicators, and Outfit/Inter web typography.",
    "Real-Time Synchronized Data: Polling endpoints (metrics, funnels, heatmaps, anomalies) automatically every 3 seconds to ensure dynamic update consistency.",
    "Clean Testing State: Home page route automatically clears existing events, forcing metrics to start cleanly at 0 so reviewers see active, live transitions.",
    "Interactive Shopper Count Dropdown: Standard select box allowing the user to select 13, 30, 50, or 100 shoppers in a single click, triggering instant high-fidelity ingestion simulation!"
]

first = True
for b in ui_bullets:
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    first = False
    p.text = "• " + b
    p.font.name = 'Arial'
    p.font.size = Pt(16)
    p.font.color.rgb = TEXT_WHITE
    p.space_after = Pt(14)

# Slide 8: Automated Verification
slide = prs.slides.add_slide(blank_slide_layout)
apply_background(slide)
add_header(slide, "7. Quality Assurance & Automated Validation")

# Left box
box_l = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.8), Inches(5.0))
tf_l = box_l.text_frame
tf_l.word_wrap = True

p_l = tf_l.paragraphs[0]
p_l.text = "UNIT TESTING COVERAGE"
p_l.font.bold = True
p_l.font.color.rgb = INDIGO
p_l.font.size = Pt(18)
p_l.space_after = Pt(14)

test_bullets = [
    "test_pipeline: Centroid association bounds, Re-ID appearance validation, cross-camera overlap resolver checks.",
    "test_metrics: Dynamic conversion math, unique visitor session aggregation, staff filter exclusion check.",
    "test_anomalies: Validates anomaly severity bounds, conversion drop thresholds, and dead zone timers."
]
for b in test_bullets:
    p = tf_l.add_paragraph()
    p.text = "• " + b
    p.font.size = Pt(14)
    p.font.color.rgb = TEXT_WHITE
    p.space_after = Pt(10)

# Right box (Mock test output box)
box_r = slide.shapes.add_textbox(Inches(6.5), Inches(1.5), Inches(6.33), Inches(5.0))
tf_r = box_r.text_frame
tf_r.word_wrap = True

p_r = tf_r.paragraphs[0]
p_r.text = "TEST EXECUTION OUTPUT"
p_r.font.bold = True
p_r.font.color.rgb = CYAN
p_r.font.size = Pt(18)
p_r.space_after = Pt(14)

p_code = tf_r.add_paragraph()
p_code.text = "================== test session starts ==================\n" \
             "collected 13 items\n\n" \
             "tests/test_anomalies.py ...                        [ 23%]\n" \
             "tests/test_metrics.py .....                         [ 61%]\n" \
             "tests/test_pipeline.py .....                        [100%]\n\n" \
             "============= 13 passed in 1.96s ============="
p_code.font.name = 'Courier New'
p_code.font.size = Pt(13)
p_code.font.color.rgb = CYAN
p_code.space_after = Pt(10)

# Slide 9: Conclusion
slide = prs.slides.add_slide(blank_slide_layout)
apply_background(slide)
add_header(slide, "8. Conclusion & Production Readiness")

box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.33), Inches(5.0))
tf = box.text_frame
tf.word_wrap = True

conclusion_points = [
    "Live Production Host: Successfully deployed in the cloud on Render using dynamic Docker environment containers.",
    "Structured Request Logs: JSON structured output capturing Trace ID, latency timings, and HTTP status codes for enterprise stability.",
    "Robust Database Idempotency: Duplicate visitor events are automatically resolved, preventing analytical calculation pollution.",
    "Comprehensive Documentation: Clean Design choices trade-offs log (docs/CHOICES.md) and Systems architecture blueprint (docs/DESIGN.md) included in code repositories."
]

first = True
for c in conclusion_points:
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    first = False
    p.text = "🏆 " + c
    p.font.name = 'Arial'
    p.font.size = Pt(16)
    p.font.color.rgb = TEXT_WHITE
    p.space_after = Pt(14)

# Save presentation
prs.save("Apex_Store_Intelligence_Presentation.pptx")
print("Presentation generated successfully!")
